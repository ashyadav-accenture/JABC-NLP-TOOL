"""Optional output: a single .pptx holding every chart as native shapes.

The PNG writers (export.py, circumplex.py, theme_priority.py) stay the
authoritative renderings -- they are what the report embeds and what gets
printed. This module re-draws the same charts as PowerPoint autoshapes so a
reader can open the deck and move a bubble, retype a label, restyle an icon or
delete a row, none of which is possible with a flattened image.

Every number is read from the same functions the PNGs use, so the deck cannot
drift from the images: the matrix comes from the matrix dataframe, the
circumplex from `resolve_layout`, and the quadrant chart from
`load_opportunities`. Nothing is re-derived here; this file is only a
different renderer.

Two deliberate differences from the PNGs:

  * Geometry is laid out in INCHES on the slide rather than in each chart's own
    data space. Matplotlib's data coordinates plus a tight bounding box have no
    meaning on a fixed 13.33 x 7.5in canvas, so each chart gets an explicit
    plot rectangle and the chart's own maths is mapped into it.

  * Text is a real text box per label, not a rendered glyph run. That is what
    makes the labels editable, and it means fine typographic details of the
    PNGs (measured bbox nudges in `_axis_triplet`, for one) are approximated by
    fixed offsets here rather than measured.
"""

from __future__ import annotations

import logging
import math
from pathlib import Path

from .models import RespondentProfile

logger = logging.getLogger(__name__)

SLIDE_W, SLIDE_H = 13.333, 7.5
INK = "1A1A1A"
GREY = "666666"
MUTED = "8A8A8A"


# ---------------------------------------------------------------------------
# Drawing helpers -- everything in inches, y measured downward from the top.
# ---------------------------------------------------------------------------

class Slide:
    """Thin drawing surface over one python-pptx slide.

    Wraps the handful of primitives the charts need (circle, rect, line,
    wedge, arc, polygon, text) so the port of each chart reads like the
    matplotlib original instead of like OOXML plumbing.
    """

    def __init__(self, prs, layout_index: int = 6):
        from pptx.util import Pt
        self.prs = prs
        self.slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        self._Pt = Pt

    # -- internals ---------------------------------------------------------

    @staticmethod
    def _emu(inches: float):
        from pptx.util import Inches
        return Inches(inches)

    @staticmethod
    def _rgb(hex_color: str):
        from pptx.dml.color import RGBColor
        return RGBColor.from_string(hex_color.lstrip("#").upper())

    def _style(self, shape, fill: str | None, line: str | None, lw: float = 1.0):
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(fill)
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = self._rgb(line)
            shape.line.width = self._Pt(lw)
        shape.shadow.inherit = False
        if shape.has_text_frame:
            shape.text_frame.text = ""
        return shape

    def _auto(self, kind, x, y, w, h, fill, line, lw=1.0):
        from pptx.enum.shapes import MSO_SHAPE
        shape = self.slide.shapes.add_shape(
            getattr(MSO_SHAPE, kind), self._emu(x), self._emu(y),
            self._emu(w), self._emu(h))
        return self._style(shape, fill, line, lw)

    # -- primitives --------------------------------------------------------

    def circle(self, cx, cy, r, fill=None, line=INK, lw=1.0, name=""):
        shape = self._auto("OVAL", cx - r, cy - r, 2 * r, 2 * r, fill, line, lw)
        if name:
            shape.name = name
        return shape

    def rect(self, x, y, w, h, fill=None, line=INK, lw=1.0, name=""):
        shape = self._auto("RECTANGLE", x, y, w, h, fill, line, lw)
        if name:
            shape.name = name
        return shape

    def line(self, x1, y1, x2, y2, color=INK, lw=1.0):
        from pptx.enum.shapes import MSO_CONNECTOR
        shape = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, self._emu(x1), self._emu(y1),
            self._emu(x2), self._emu(y2))
        shape.line.color.rgb = self._rgb(color)
        shape.line.width = self._Pt(lw)
        return shape

    # Curved geometry is emitted as freeform point lists rather than as the
    # ARC / PIE / BLOCK_ARC preset shapes. The presets carry their sweep in
    # adjustment values, and renderers disagree badly about them -- macOS
    # QuickLook draws nothing at all for a default PIE -- which would silently
    # drop every persona ring and every smiley mouth on some machines. A
    # freeform is drawn by everything, and PowerPoint still exposes its points
    # for editing.
    @staticmethod
    def _arc_points(cx, cy, rx, ry, theta1, theta2, segments=48):
        """Points along a matplotlib-convention arc (counter-clockwise from 3
        o'clock), converted to slide space where y points down."""
        if theta2 < theta1:
            theta2 += 360.0
        step = (theta2 - theta1) / segments
        return [(cx + rx * math.cos(math.radians(theta1 + i * step)),
                 cy - ry * math.sin(math.radians(theta1 + i * step)))
                for i in range(segments + 1)]

    def wedge(self, cx, cy, r, theta1, theta2, width, fill):
        """Donut segment: outer arc out, inner arc back."""
        outer = self._arc_points(cx, cy, r, r, theta1, theta2)
        inner = self._arc_points(cx, cy, r - width, r - width, theta1, theta2)
        return self.poly(outer + inner[::-1], fill=fill, line=None)

    def pie(self, cx, cy, r, theta1, theta2, fill):
        return self.poly([(cx, cy)] + self._arc_points(cx, cy, r, r, theta1, theta2),
                         fill=fill, line=None)

    def arc(self, cx, cy, rx, ry, theta1, theta2, color=INK, lw=1.0):
        """Open arc (stroke only), matching matplotlib's Arc(width, height)."""
        return self.poly(self._arc_points(cx, cy, rx, ry, theta1, theta2),
                         fill=None, line=color, lw=lw, close=False)

    def poly(self, points, fill=None, line=INK, lw=1.0, close=True):
        """Freeform from (x, y) inch pairs."""
        builder = self.slide.shapes.build_freeform(
            self._emu(points[0][0]), self._emu(points[0][1]))
        builder.add_line_segments(
            [(self._emu(px), self._emu(py)) for px, py in points[1:]], close=close)
        return self._style(builder.convert_to_shape(), fill, line, lw)

    def text(self, x, y, s, size=10.0, bold=False, italic=False, color=INK,
             ha="center", va="center", width=2.0, rotation=0.0, line_spacing=None):
        """A text box whose (x, y) is the anchor point named by ha/va.

        Boxes are sized rather than auto-fit: PowerPoint only recomputes
        autofit when a human edits the text, so an unsized box would render
        at whatever height the writer guessed anyway.
        """
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        # A point is 1/72in; the 1.22 is the line box around the glyphs. The
        # height only has to be right enough for ha/va anchoring to land --
        # PowerPoint overflows rather than clips if a run wraps further.
        n_lines = s.count("\n") + 1
        height = size / 72.0 * 1.22 * (line_spacing or 1.0) * n_lines + 0.06
        left = {"center": x - width / 2, "left": x, "right": x - width}[ha]
        top = {"center": y - height / 2, "top": y, "bottom": y - height}[va]
        box = self.slide.shapes.add_textbox(
            self._emu(left), self._emu(top), self._emu(width), self._emu(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, chunk in enumerate(s.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT,
                              "right": PP_ALIGN.RIGHT}[ha]
            if line_spacing:
                para.line_spacing = line_spacing
            run = para.add_run()
            run.text = chunk
            run.font.size = self._Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = self._rgb(color)
        if rotation:
            box.rotation = rotation
        return box


# ---------------------------------------------------------------------------
# Icons -- ports of circumplex.py's matplotlib drawers onto Slide
# ---------------------------------------------------------------------------
# Same proportions as the originals: each fills a box of half-width `s`
# centred on (cx, cy). Only the y sign flips, because slide y points down.

def _icon_book(sl, cx, cy, s, c):
    for sx in (-1, 1):
        sl.poly([(cx, cy - s * 0.52), (cx + sx * s * 0.72, cy - s * 0.34),
                 (cx + sx * s * 0.72, cy + s * 0.46), (cx, cy + s * 0.30)],
                fill=None, line=c, lw=1.1)
    sl.line(cx, cy - s * 0.52, cx, cy + s * 0.30, c, 1.1)


def _icon_package(sl, cx, cy, s, c):
    sl.rect(cx - s * 0.62, cy - s * 0.50, s * 1.24, s * 1.0, None, c, 1.1)
    sl.line(cx - s * 0.62, cy - s * 0.14, cx + s * 0.62, cy - s * 0.14, c, 1.0)
    sl.line(cx, cy - s * 0.50, cx, cy - s * 0.14, c, 1.0)


def _icon_referral(sl, cx, cy, s, c):
    for dx, dy, w, h in ((-0.30, -0.18, 0.80, 0.46), (0.26, 0.20, 0.70, 0.40)):
        sl.rect(cx + s * dx - s * w / 2, cy + s * dy - s * h / 2,
                s * w, s * h, "FFFFFF", c, 1.0)
    sl.poly([(cx - s * 0.44, cy + s * 0.05), (cx - s * 0.30, cy + s * 0.26),
             (cx - s * 0.18, cy + s * 0.05)], fill=c, line=None)


def _icon_spark(sl, cx, cy, s, c):
    pts = []
    for i in range(8):
        ang = math.radians(90 + i * 45)
        r = s * (0.66 if i % 2 == 0 else 0.24)
        pts.append((cx + r * math.cos(ang), cy - r * math.sin(ang)))
    sl.poly(pts, fill=None, line=c, lw=1.1)


def _icon_expert(sl, cx, cy, s, c):
    sl.circle(cx, cy - s * 0.30, s * 0.28, None, c, 1.1)
    sl.arc(cx, cy + s * 0.42, s * 0.58, s * 0.50, 20, 160, c, 1.1)


def _icon_stale_doc(sl, cx, cy, s, c):
    sl.poly([(cx - s * 0.48, cy - s * 0.56), (cx + s * 0.22, cy - s * 0.56),
             (cx + s * 0.50, cy - s * 0.26), (cx + s * 0.50, cy + s * 0.56),
             (cx - s * 0.48, cy + s * 0.56)], fill=None, line=c, lw=1.1)
    for k, dy in enumerate((-0.06, 0.16, 0.38)):
        sl.line(cx - s * 0.28, cy + s * dy, cx + s * (0.28 - k * 0.10), cy + s * dy, c, 0.9)


def _icon_hourglass(sl, cx, cy, s, c):
    sl.poly([(cx - s * 0.46, cy - s * 0.56), (cx + s * 0.46, cy - s * 0.56), (cx, cy)],
            fill=c, line=c, lw=1.0)
    sl.poly([(cx - s * 0.46, cy + s * 0.56), (cx + s * 0.46, cy + s * 0.56), (cx, cy)],
            fill=None, line=c, lw=1.0)
    for dy in (-0.60, 0.60):
        sl.line(cx - s * 0.54, cy + s * dy, cx + s * 0.54, cy + s * dy, c, 1.2)


def _icon_calendar(sl, cx, cy, s, c):
    sl.rect(cx - s * 0.56, cy - s * 0.48, s * 1.12, s * 1.0, None, c, 1.1)
    sl.rect(cx - s * 0.56, cy - s * 0.48, s * 1.12, s * 0.24, c, c, 0.8)
    for gx in (-0.28, 0.0, 0.28):
        for gy in (0.04, 0.30):
            sl.rect(cx + s * gx - 0.012, cy + s * gy - 0.012, 0.024, 0.024, c, None)


def _icon_monitor(sl, cx, cy, s, c):
    sl.rect(cx - s * 0.62, cy - s * 0.60, s * 1.24, s * 0.80, None, c, 1.1)
    sl.line(cx - s * 0.62, cy - s * 0.38, cx + s * 0.62, cy - s * 0.38, c, 0.9)
    sl.line(cx, cy + s * 0.20, cx, cy + s * 0.44, c, 1.1)
    sl.line(cx - s * 0.34, cy + s * 0.44, cx + s * 0.34, cy + s * 0.44, c, 1.2)


def _icon_warning(sl, cx, cy, s, c):
    sl.poly([(cx, cy - s * 0.58), (cx + s * 0.62, cy + s * 0.48),
             (cx - s * 0.62, cy + s * 0.48)], fill=None, line=c, lw=1.1)
    sl.line(cx, cy - s * 0.24, cx, cy + s * 0.14, c, 1.3)
    sl.circle(cx, cy + s * 0.31, s * 0.07, c, None)


def _icon_puzzle(sl, cx, cy, s, c):
    sl.rect(cx - s * 0.52, cy - s * 0.52, s * 1.04, s * 1.04, None, c, 1.1)
    sl.circle(cx + s * 0.52, cy, s * 0.20, "FFFFFF", c, 1.1)
    sl.circle(cx, cy - s * 0.52, s * 0.20, "FFFFFF", c, 1.1)


def _icon_megaphone(sl, cx, cy, s, c):
    sl.poly([(cx - s * 0.62, cy - s * 0.30), (cx - s * 0.62, cy + s * 0.30),
             (cx + s * 0.06, cy + s * 0.58), (cx + s * 0.06, cy - s * 0.58)],
            fill=None, line=c, lw=1.1)
    for scale in (0.44, 0.78):
        sl.arc(cx + s * 0.10, cy, s * scale / 2, s * scale * 0.75, -58, 58, c, 1.0)


ICON_DRAWERS = {
    "book": _icon_book, "package": _icon_package, "referral": _icon_referral,
    "spark": _icon_spark, "expert": _icon_expert, "stale_doc": _icon_stale_doc,
    "hourglass": _icon_hourglass, "calendar": _icon_calendar,
    "monitor": _icon_monitor, "warning": _icon_warning,
    "puzzle": _icon_puzzle, "megaphone": _icon_megaphone,
}


def _draw_icon(sl, name, cx, cy, size, color):
    drawer = ICON_DRAWERS.get(name)
    if drawer is None:
        logger.warning("Unknown circumplex icon '%s'; drawing nothing.", name)
        return
    drawer(sl, cx, cy, size, color)


# ---------------------------------------------------------------------------
# Score faces -- port of export.py's _draw_score_face
# ---------------------------------------------------------------------------

def _draw_face(sl, cx, cy, tier, size=0.16):
    """Same construction as the PNG's faces: hue, mouth shape, mouth weight
    and brow treatment all move together, so the ranking survives greyscale."""
    face = tier["face"]
    curve = tier["curve"]
    sl.circle(cx, cy, size, face, None)

    eye_dx, eye_dy, eye_r = size * 0.36, size * 0.28, size * 0.10
    if curve >= 0.9:
        for sx in (-1, 1):
            sl.arc(cx + sx * eye_dx, cy - eye_dy, size * 0.20, size * 0.17,
                   20, 160, "000000", 1.4)
    else:
        for sx in (-1, 1):
            sl.circle(cx + sx * eye_dx, cy - eye_dy, eye_r, "000000", None)
        if curve <= -0.9:
            for sx in (-1, 1):
                sl.line(cx + sx * size * 0.56, cy - size * 0.44,
                        cx + sx * size * 0.20, cy - size * 0.62, "000000", 1.4)

    mouth_w = size * 1.0
    if abs(curve) < 0.05:
        sl.line(cx - mouth_w / 2, cy + size * 0.34, cx + mouth_w / 2,
                cy + size * 0.34, "000000", 1.8)
    elif curve >= 0.9:
        sl.pie(cx, cy + size * 0.04, size * 0.62, 200, 340, "000000")
    else:
        mouth_h = max(size * 0.28, size * 0.62 * abs(curve))
        if curve > 0:
            center_y = cy + size * 0.46 - mouth_h / 2
            theta1, theta2 = 200, 340
        else:
            center_y = cy + size * 0.24 + mouth_h / 2
            theta1, theta2 = 20, 160
        sl.arc(cx, center_y, mouth_w / 2, mouth_h / 2, theta1, theta2, "000000", 1.8)


# ---------------------------------------------------------------------------
# Slide 1 -- Brand Health Matrix + score scale key
# ---------------------------------------------------------------------------

def _slide_matrix(prs, matrix_df, config) -> None:
    from .export import HEADER_PURPLE, SMILEY_TIERS, _safe_mean, _smiley_tier
    from .matrix import suppressed_cells

    import pandas as pd

    sl = Slide(prs)
    themes = ["Awareness", "Understanding", "Trust", "Relevance",
              "Ease of Engagement", "Advocacy"]
    personas = matrix_df["Brand Persona"].tolist()
    blanked = suppressed_cells(config)

    sl.text(0.55, 0.34, "JABC Brand Health Matrix", size=26, bold=True,
            ha="left", va="top", width=8.0)

    left, top = 0.55, 1.10
    label_w, col_w = 1.75, 1.53
    head_h, row_h = 0.62, 0.72
    n_cols = len(themes) + 1

    sl.rect(left, top, label_w + n_cols * col_w, head_h, HEADER_PURPLE, None)
    sl.text(left + label_w / 2, top + head_h / 2, "Brand\nPersona", size=10.5,
            bold=True, color="FFFFFF", width=label_w)
    for j, theme in enumerate(themes):
        sl.text(left + label_w + (j + 0.5) * col_w, top + head_h / 2, theme,
                size=10, bold=True, color="FFFFFF", width=col_w - 0.06)
    sl.text(left + label_w + (len(themes) + 0.5) * col_w, top + head_h / 2,
            "Average\nof Personas", size=9.5, bold=True, italic=True,
            color="FFFFFF", width=col_w - 0.06)

    def cell(x, y, value, suppressed=False, numeric=False):
        if suppressed:
            sl.rect(x, y, col_w, row_h, "FFFFFF", "BBBBBB", 0.75)
            return
        tier = _smiley_tier(value)
        sl.rect(x, y, col_w, row_h, tier["bg"] if tier else "FFFFFF", "BBBBBB", 0.75)
        has_value = value is not None and not (isinstance(value, float) and pd.isna(value))
        if numeric:
            sl.text(x + col_w / 2, y + row_h / 2, f"{value:.2f}" if has_value else "–",
                    size=12, bold=True, italic=True, color="333333", width=col_w)
        elif tier is None:
            sl.text(x + col_w / 2, y + row_h / 2, "–", size=13, color="999999", width=col_w)
        else:
            _draw_face(sl, x + col_w / 2, y + row_h / 2, tier, size=row_h * 0.30)

    for i, persona in enumerate(personas):
        y = top + head_h + i * row_h
        sl.rect(left, y, label_w, row_h, "FFFFFF", "BBBBBB", 0.75)
        sl.text(left + label_w / 2, y + row_h / 2, persona, size=10, bold=True,
                width=label_w - 0.08)
        for j, theme in enumerate(themes):
            cell(left + label_w + j * col_w, y, matrix_df.iloc[i][theme],
                 suppressed=(persona, theme) in blanked)
        cell(left + label_w + len(themes) * col_w, y,
             matrix_df.iloc[i]["Average Score"], numeric=True)

    y = top + head_h + len(personas) * row_h
    sl.rect(left, y, label_w, row_h, HEADER_PURPLE, "BBBBBB", 0.75)
    sl.text(left + label_w / 2, y + row_h / 2, "Average\nof Themes", size=10,
            bold=True, italic=True, color="FFFFFF", width=label_w - 0.08)
    theme_averages = [_safe_mean(matrix_df[t]) for t in themes]
    for j, avg in enumerate(theme_averages):
        cell(left + label_w + j * col_w, y, avg, numeric=True)
    cell(left + label_w + len(themes) * col_w, y, _safe_mean(theme_averages), numeric=True)

    # --- Score scale key, the same faces the cells use ---------------------
    key_top = y + row_h + 0.72
    total_w = label_w + n_cols * col_w
    swatch_w = total_w / len(SMILEY_TIERS)
    sl.text(left, key_top - 0.62, "Brand Health Score Scale", size=12, bold=True,
            ha="left", va="top", width=4.0)
    sl.text(left, key_top - 0.26, "worse", size=9, italic=True, color="999999",
            ha="left", va="top", width=0.6)
    sl.text(left + total_w, key_top - 0.26, "better", size=9, italic=True,
            color="999999", ha="right", va="top", width=0.6)
    # The direction rail: which way is good is stated, not left to colour
    # convention, exactly as in the standalone PNG key.
    sl.line(left + 0.55, key_top - 0.18, left + total_w - 0.55, key_top - 0.18,
            "999999", 1.0)

    for k, tier in enumerate(reversed(SMILEY_TIERS)):
        x0 = left + k * swatch_w
        sl.rect(x0, key_top, swatch_w - 0.10, 0.62, tier["bg"], "BBBBBB", 0.75)
        _draw_face(sl, x0 + 0.28, key_top + 0.31, tier, size=0.17)
        lo = tier["min"]
        higher = SMILEY_TIERS[len(SMILEY_TIERS) - 1 - k - 1]["min"] if k < len(SMILEY_TIERS) - 1 else None
        if lo == float("-inf"):
            band = f"below {higher:.1f}" if higher is not None else "lowest"
        elif higher is None:
            band = f"{lo:.1f} and up"
        else:
            band = f"{lo:.1f} – {higher:.1f}"
        sl.text(x0 + 0.52, key_top + 0.19, tier["label"], size=10.5, bold=True,
                color="333333", ha="left", width=swatch_w - 0.6)
        sl.text(x0 + 0.52, key_top + 0.43, band, size=9, color=GREY,
                ha="left", width=swatch_w - 0.6)


# ---------------------------------------------------------------------------
# Slide 2 -- Factor Circumplex
# ---------------------------------------------------------------------------

def _slide_circumplex(prs, profiles, config) -> bool:
    from .circumplex import (KIND_RING_COLORS, KIND_STYLE, RING_W, Y_SPAN,
                             _dot_radius, resolve_layout)

    records = resolve_layout(profiles, config)
    if not records:
        logger.warning("No circumplex_layout configured; skipping circumplex slide.")
        return False

    raw = getattr(config, "motivators_barriers", {}) or {}
    layout = raw.get("circumplex_layout", {}) or {}
    axis_labels = layout.get("axis_labels", {})
    quadrants = layout.get("quadrant_labels", {})

    sl = Slide(prs)
    sl.text(0.55, 0.30, "Factor Circumplex", size=26, bold=True, ha="left",
            va="top", width=8.0)

    # The chart's data circle has radius R=5; RIN is that radius on the slide.
    # cy clears the title, and cx leaves room for the left axis label, which
    # is the longest piece of furniture on the slide.
    R, RIN = 5.0, 2.42
    cx, cy = 4.55, 4.20
    scale = RIN / R

    def px(x, y):
        return cx + x * scale, cy - y * scale

    sl.circle(cx, cy, RIN, None, "2F5C99", 1.3)
    sl.line(cx - RIN, cy, cx + RIN, cy, "333333", 0.9)
    sl.line(cx, cy - RIN, cx, cy + RIN, "333333", 0.9)

    sl.text(cx, cy - RIN * 1.10, axis_labels.get("top", "HIGH POSITIVE EFFECT"),
            size=9, bold=True, color="333333", va="bottom", width=3.2)
    note = axis_labels.get("top_note", "")
    if note:
        sl.text(cx, cy - RIN * 1.045, f"({note})", size=8, italic=True,
                color=MUTED, va="bottom", width=2.4)
    sl.text(cx, cy + RIN * 1.10, axis_labels.get("bottom", "HIGH NEGATIVE EFFECT"),
            size=9, bold=True, color="333333", va="top", width=3.2)
    sl.text(cx - RIN * 1.06, cy, axis_labels.get("left", ""), size=9, bold=True,
            color="333333", ha="right", width=1.5)
    sl.text(cx + RIN * 1.06, cy, axis_labels.get("right", ""), size=9, bold=True,
            color="333333", ha="left", width=1.6)

    for angle_deg, key in ((135, "upper_left"), (45, "upper_right"),
                           (225, "lower_left"), (315, "lower_right")):
        text = quadrants.get(key, "")
        if not text:
            continue
        rad = math.radians(angle_deg)
        lx, ly = px(R * 1.22 * math.cos(rad), R * 1.22 * math.sin(rad))
        # Slide rotation is clockwise, matplotlib's is counter-clockwise.
        rotation = -(((angle_deg - 90 + 90) % 180) - 90)
        sl.text(lx, ly, text, size=8.5, bold=True, color=MUTED, width=1.9,
                rotation=rotation)

    frequencies = [r["frequency"] for r in records]
    min_freq, max_freq = min(frequencies), max(frequencies)
    for rec in records:
        r = _dot_radius(rec["frequency"], min_freq, max_freq)
        sign = 1 if rec["kind"] == "motivator" else -1
        x, y = px(rec["x"] * R, sign * (rec["frequency"] / max_freq) * R * Y_SPAN)
        rec["_pos"] = (x, y, r * scale)

    # Same three passes as the PNG: every fill, then every icon, then every
    # ring. Circles overlap (B3 and B4 tie on frequency), and drawing each one
    # complete in turn would let a later fill cover an earlier ring. The ring is
    # a solid family band now; the persona breakdown it used to carry has its
    # own slides, one pie per factor.
    for rec in records:
        x, y, r = rec["_pos"]
        style = KIND_STYLE[rec["kind"]]
        sl.circle(x, y, r, style["fill"].lstrip("#"), style["edge"].lstrip("#"), 1.1,
                  name=f"{rec['code']} disc")
    for rec in records:
        x, y, r = rec["_pos"]
        _draw_icon(sl, rec["icon"], x, y, r * 0.62, KIND_STYLE[rec["kind"]]["ink"].lstrip("#"))
    for rec in records:
        x, y, r = rec["_pos"]
        # Two half-sweeps rather than one 0-360 wedge: a freeform whose start
        # and end points coincide closes on itself and some renderers drop it.
        for start in (0.0, 180.0):
            sl.wedge(x, y, r, start, start + 180.0, RING_W * scale,
                     KIND_RING_COLORS[rec["kind"]].lstrip("#"))
    for rec in records:
        x, y, r = rec["_pos"]
        sl.text(x, y + r + 0.04, rec["code"], size=7.5, bold=True,
                color=KIND_STYLE[rec["kind"]]["ink"].lstrip("#"), va="top", width=0.7)

    # --- Key, to the right of the plot rather than beneath it: the slide is
    # wide where the PNG is tall, so a stacked key would waste the right third.
    kx = 7.65
    ky = 1.05
    sl.text(kx, ky, "Factors", size=12, bold=True, color="333333", ha="left",
            va="top", width=1.2)
    sl.text(kx + 0.95, ky + 0.05, "(coded frequency of mention)", size=8.5,
            italic=True, color=MUTED, ha="left", va="top", width=2.6)
    row_h = 0.33
    for i, rec in enumerate(records):
        row_y = ky + 0.45 + i * row_h
        style = KIND_STYLE[rec["kind"]]
        sl.circle(kx + 0.12, row_y + row_h / 2, 0.11, style["fill"].lstrip("#"),
                  style["edge"].lstrip("#"), 0.9)
        _draw_icon(sl, rec["icon"], kx + 0.12, row_y + row_h / 2, 0.075,
                   style["ink"].lstrip("#"))
        sl.text(kx + 0.32, row_y + row_h / 2,
                f"{rec['code']}  {rec['label']}  ({rec['frequency']})",
                size=8.5, color="333333", ha="left", width=5.3)

    py = ky + 0.45 + len(records) * row_h + 0.28
    sl.text(kx, py, "Ring", size=12, bold=True, color="333333",
            ha="left", va="top", width=1.6)
    sl.text(kx + 0.45, py + 0.05, "(which family the factor belongs to)", size=8.5,
            italic=True, color=MUTED, ha="left", va="top", width=3.0)
    for i, (kind, label) in enumerate((("motivator", "Motivator"), ("barrier", "Barrier"))):
        x0 = kx + i * 2.55
        row_y = py + 0.45
        sl.circle(x0 + 0.10, row_y, 0.085, KIND_RING_COLORS[kind].lstrip("#"), None)
        sl.text(x0 + 0.28, row_y, label, size=8.5, color="333333", ha="left", width=2.2)

    sl.text(kx, py + 0.45 + row_h + 0.30,
            "Vertical position and circle size = coded frequency of mention;\n"
            "motivators read upward, barriers downward. Horizontal position is\n"
            "set by the approved diagram — an editorial judgement, not a\n"
            "measurement, so it is never recomputed from the data.\n"
            "Who raised each factor follows, one pie per factor.",
            size=8.5, italic=True, color=MUTED, ha="left", va="top", width=5.4,
            line_spacing=1.35)
    return True


# ---------------------------------------------------------------------------
# Slides 3-4 -- who raised each factor, one pie per factor
# ---------------------------------------------------------------------------
# These carry the persona breakdown the circumplex rings used to hold. At ring
# size a two-slice donut on ten circles, several of them overlapping, could be
# seen but not read -- no slice was labelled and no percentage was recoverable.
# Full-size pies with their own titles and printed percentages say the same
# thing legibly, and the circumplex gets a clean family ring in exchange.

def _slide_factor_pies(prs, records, kind: str, heading: str, config) -> bool:
    from .circumplex import KIND_RING_COLORS, group_color, group_shares
    from .config_loader import persona_group_display_names, persona_groups

    items = [r for r in records if r["kind"] == kind]
    if not items:
        return False

    group_names = persona_group_display_names(config)
    accent = KIND_RING_COLORS[kind].lstrip("#")

    sl = Slide(prs)
    sl.text(0.55, 0.30, heading, size=26, bold=True, ha="left", va="top", width=9.0)
    sl.text(0.55, 0.95, "Share of coded mentions by Brand Persona. Percentages are of "
                        "that factor's own mentions, so each pie totals 100%.",
            size=10.5, italic=True, color=MUTED, ha="left", va="top", width=11.0)

    # One column per factor across the full slide width. Five factors per family
    # today; the arithmetic holds for any count that still leaves a legible pie.
    left, right = 0.55, SLIDE_W - 0.55
    col_w = (right - left) / len(items)
    r = min(1.05, col_w * 0.40)
    cy = 3.35

    for i, rec in enumerate(items):
        cx = left + col_w * (i + 0.5)

        sl.text(cx, 1.55, f"{rec['code']}  {rec['label']}", size=10.5, bold=True,
                color="333333", va="top", width=col_w - 0.18)

        shares = {k: v for k, v in group_shares(rec["persona_shares"], config).items()
                  if v > 0}
        total = sum(shares.values())
        # analyze_drivers substitutes an even split when no cue fired, which is
        # readable as "unknown" on a small ring but not on a labelled pie -- an
        # even split printed as 50% / 50% is indistinguishable from a measured
        # one. So the unattributed case is drawn as an empty circle and says so.
        if rec["mention_count"] <= 0 or total <= 0:
            sl.circle(cx, cy, r, "EFEFEF", "CFCFCF", 1.0)
            sl.text(cx, cy, "no attributable\nmentions", size=9.5, italic=True,
                    color=MUTED, width=col_w - 0.30)
        else:
            start = 90.0
            for gi, (group_key, share) in enumerate(shares.items()):
                # Renormalised so rounding in the shares cannot leave a sliver of
                # background showing through the pie.
                sweep = share / total * 360.0
                fill = group_color(group_key, gi).lstrip("#")
                sl.pie(cx, cy, r, start, start + sweep, fill)
                mid = math.radians(start + sweep / 2)
                # Percentages sit inside the wedge, except on a slice too thin to
                # hold the text, which gets a line out to the side instead.
                if sweep >= 42:
                    sl.text(cx + 0.58 * r * math.cos(mid), cy - 0.58 * r * math.sin(mid),
                            f"{share / total * 100:.0f}%", size=12, bold=True,
                            color="FFFFFF", width=0.9)
                else:
                    lx = cx + 1.22 * r * math.cos(mid)
                    ly = cy - 1.22 * r * math.sin(mid)
                    sl.line(cx + 0.95 * r * math.cos(mid), cy - 0.95 * r * math.sin(mid),
                            lx, ly, "8A8A8A", 0.75)
                    sl.text(lx, ly, f"{share / total * 100:.0f}%", size=10, bold=True,
                            color="333333", width=0.7)
                start += sweep

        if rec["mention_count"] <= 0:
            detail = (f"{rec['frequency']} coded mentions\n"
                      "no cue matched in the transcripts")
        else:
            detail = (f"{rec['frequency']} coded mentions\n"
                      f"split from {rec['mention_count']} detected mentions "
                      f"across {rec['respondent_count']} respondents")
            if not rec["well_evidenced"]:
                detail += "\n(thin evidence — read with care)"
        sl.text(cx, cy + r + 0.28, detail, size=8.5, color=MUTED, va="top",
                width=col_w - 0.18, line_spacing=1.25)

    ly = 6.55
    sl.rect(left, ly - 0.06, right - left, 0.02, accent, None)
    for i, group_key in enumerate(persona_groups(config)):
        x0 = left + i * 3.2
        sl.circle(x0 + 0.10, ly + 0.32, 0.10, group_color(group_key, i).lstrip("#"), None)
        sl.text(x0 + 0.30, ly + 0.32, group_names.get(group_key, group_key), size=10,
                color="333333", ha="left", width=2.7)
    sl.text(right, ly + 0.32,
            "Attribution comes from the interview text; the coded frequency does not.",
            size=8.5, italic=True, color=MUTED, ha="right", width=5.2)
    return True


# ---------------------------------------------------------------------------
# Slide 5 -- Opportunity priority matrix
# ---------------------------------------------------------------------------

def _slide_opportunities(prs, profiles, config) -> bool:
    from .opportunities import load_opportunities
    from .theme_priority import (BOX_H, BOX_W, BUBBLE_PURPLE, EYEBROW,
                                 FAINT_LAVENDER, PRIORITISE_ACCENT,
                                 PRIORITISE_FILL, QUAD_LABEL_PURPLE, STRETCH_TO_FILL,
                                 SUBTITLE, SUBTITLE_GREY, TITLE, _bubble_radius,
                                 _declutter, _opportunity_metrics, _stretch_axis, _to_box)

    opportunities = load_opportunities(config)
    if not opportunities:
        logger.warning("No opportunities configured; skipping opportunity slide.")
        return False

    sl = Slide(prs)
    sl.text(0.55, 0.28, " ".join(EYEBROW.upper()), size=7.5, bold=True,
            color="7A2E9F", ha="left", va="top", width=4.0)
    sl.text(0.55, 0.44, TITLE, size=26, bold=True, ha="left", va="top", width=8.0)
    sl.text(0.55, 1.00, SUBTITLE, size=9.5, color=SUBTITLE_GREY.lstrip("#"),
            ha="left", va="top", width=11.0, line_spacing=1.35)

    # Plot box on the slide, and the same value/difficulty maths as the PNG
    # mapped into it.
    bx, by, bw, bh = 1.55, 1.85, 7.95, 4.15
    items = [o.code for o in opportunities]
    by_code = {o.code: o for o in opportunities}
    metrics = {o.code: _opportunity_metrics(o) for o in opportunities}
    max_value = max(m[0] for m in metrics.values())
    max_count = max((m[2] for m in metrics.values()), default=0)
    # Radii are computed in the PNG's data units, so they scale with the box.
    radii = {c: _bubble_radius(m[2], max_count) for c, m in metrics.items()}
    positions = {c: _to_box(m[0], m[1], max_value) for c, m in metrics.items()}
    if STRETCH_TO_FILL:
        xs = _stretch_axis({t: p[0] for t, p in positions.items()}, BOX_W)
        ys = _stretch_axis({t: p[1] for t, p in positions.items()}, BOX_H)
        positions = {t: (xs[t], ys[t]) for t in positions}
    # No keep-out zones: the corner labels sit in the box's corners on the
    # slide too, so the PNG's measured text bboxes are approximated by fixed
    # corner rectangles in data units.
    pad = 3.0
    zones = [(0, BOX_H - 12, 26, BOX_H), (BOX_W - 30, BOX_H - 12, BOX_W, BOX_H),
             (0, 0, 34, 12), (BOX_W - 30, 0, BOX_W, 12)]
    positions = _declutter(positions, radii, zones)

    sx, sy = bw / BOX_W, bh / BOX_H
    # Bubbles must stay round, so they take the smaller of the two scales.
    rs = min(sx, sy)

    def px(x, y):
        return bx + x * sx, by + bh - y * sy

    sl.rect(bx, by, bw, bh, None, FAINT_LAVENDER.lstrip("#"), 1.1)
    sl.line(bx + bw / 2, by, bx + bw / 2, by + bh, FAINT_LAVENDER.lstrip("#"), 1.1)
    sl.line(bx, by + bh / 2, bx + bw, by + bh / 2, FAINT_LAVENDER.lstrip("#"), 1.1)
    sl.rect(bx + bw / 2, by + bh / 2, bw / 2, bh / 2, PRIORITISE_FILL.lstrip("#"),
            PRIORITISE_ACCENT.lstrip("#"), 1.8)

    corner = [("Avoid", pad, BOX_H - pad, "left", "top", FAINT_LAVENDER),
              ("Consider", BOX_W - pad, BOX_H - pad, "right", "top", QUAD_LABEL_PURPLE),
              ("Investigate", pad, pad, "left", "bottom", QUAD_LABEL_PURPLE),
              ("Prioritise", BOX_W - pad, pad, "right", "bottom", PRIORITISE_ACCENT)]
    for label, lx, ly, ha, va, color in corner:
        tx, ty = px(lx, ly)
        sl.text(tx, ty, label, size=14, bold=True, color=color.lstrip("#"),
                ha=ha, va=va, width=1.6)

    # Axis captions: "Low — <axis> — High", mirrored on both sides as in the
    # source deck. Three text boxes, so only the axis name is bold.
    def axis_triplet(x, y, name, rotation=0.0):
        gap = 0.66 if rotation else 0.78
        sl.text(x, y, name, size=10.5, bold=True, width=2.1, rotation=rotation)
        if rotation:
            sl.text(x, y + gap, "Low –", size=10.5, width=0.7, rotation=rotation)
            sl.text(x, y - gap, "– High", size=10.5, width=0.7, rotation=rotation)
        else:
            sl.text(x - gap, y, "Low–", size=10.5, ha="right", width=0.7)
            sl.text(x + gap, y, "– High", size=10.5, ha="left", width=0.7)

    axis_triplet(bx + bw / 2, by - 0.28, "Relative Value")
    axis_triplet(bx + bw / 2, by + bh + 0.30, "Relative Value")
    sl.text(bx + bw / 2, by + bh + 0.52, "number of educators who raised the theme",
            size=8.5, color=SUBTITLE_GREY.lstrip("#"), va="top", width=4.0)
    # Rotated about their own centres, so the boxes stay inside the margins
    # the plot box leaves free on each side.
    axis_triplet(bx - 0.70, by + bh / 2, "Relative Difficulty", rotation=-90)
    axis_triplet(bx + bw + 0.70, by + bh / 2, "Relative Difficulty", rotation=90)

    for code in items:
        x, y = px(*positions[code])
        r = radii[code] * rs
        opp = by_code[code]
        sl.circle(x, y, r, BUBBLE_PURPLE.lstrip("#"), None, name=f"{code} bubble")
        sl.text(x, y - r * 0.16, code, size=max(9.0, min(13.0, r * 26)), bold=True,
                color="FFFFFF", width=0.8)
        sl.text(x, y + r * 0.42, f"{opp.mentions}", size=7.5, color="FFFFFF", width=0.5)

    # Key: two columns beside the box, where the PNG stacks them underneath.
    kx, ky = 10.45, 1.95
    for i, code in enumerate(items):
        opp = by_code[code]
        row_y = ky + i * 0.42
        sl.text(kx, row_y, code, size=9, bold=True, color=BUBBLE_PURPLE.lstrip("#"),
                ha="left", va="top", width=0.4)
        sl.text(kx + 0.34, row_y, f"{opp.title}  ({opp.mentions})", size=8,
                color="333333", ha="left", va="top", width=2.4, line_spacing=1.2)

    legend_y = ky + len(items) * 0.42 + 0.25
    sl.circle(kx + 0.1, legend_y + 0.09, 0.1, BUBBLE_PURPLE.lstrip("#"), None)
    sl.text(kx + 0.32, legend_y, "Size = number of people who raised it", size=8.5,
            bold=True, color="333333", ha="left", va="top", width=2.6)
    sl.text(kx, legend_y + 0.34,
            "Value = coded mentions (measured).\nDifficulty = authored estimate\n"
            "(config/opportunities.yaml).",
            size=7.5, italic=True, color=MUTED, ha="left", va="top", width=2.8,
            line_spacing=1.3)
    return True


# ---------------------------------------------------------------------------
# Deck
# ---------------------------------------------------------------------------

def write_chart_deck(profiles: list[RespondentProfile], matrix_df, config,
                     path: Path) -> bool:
    """Writes every chart into one .pptx as editable shapes.

    Returns False if python-pptx isn't installed, mirroring the optional PNG
    writers' behaviour when matplotlib is missing.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        logger.warning("python-pptx not available; skipping chart deck.")
        return False

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    from .circumplex import resolve_layout

    _slide_matrix(prs, matrix_df, config)
    _slide_circumplex(prs, profiles, config)
    # resolve_layout is cheap and deterministic, so the pie slides re-read it
    # rather than taking the circumplex slide's records: those carry `_pos`
    # keys in slide inches, which would be meaningless here.
    records = resolve_layout(profiles, config)
    _slide_factor_pies(prs, records, "motivator",
                       "Motivators — who raised each factor", config)
    _slide_factor_pies(prs, records, "barrier",
                       "Barriers — who raised each factor", config)
    _slide_opportunities(prs, profiles, config)

    prs.save(str(path))
    return True
